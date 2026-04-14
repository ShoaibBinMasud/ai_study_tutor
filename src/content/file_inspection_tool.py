"""
FileInspectionTool — Simple file metadata extraction.

Inspects files and returns their metadata:
- file type, page count, is_scanned, has_toc, toc_entry_count

Handles:
- Single file path
- Multiple file paths (list)
- Folder path (scans all supported files)

Zero LLM calls. Pure Python.
"""

import logging
from pathlib import Path
from typing import Optional, Dict, Any, List, Union

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".pptx", ".docx",
    ".png", ".jpg", ".jpeg", ".webp",
    ".txt", ".md"
}


def inspect(path_input: Union[str, List[str]]) -> List[Dict[str, Any]]:
    """
    Inspect files intelligently.

    Args:
        path_input: Can be:
            - Single file path (str)
            - List of file paths
            - Folder path (will scan all supported files)

    Returns:
        List of metadata dicts for all files found.
    """
    file_paths = []

    # Normalize input to list of file paths
    if isinstance(path_input, str):
        path = Path(path_input)
        if path.is_dir():
            # Folder: find all supported files
            file_paths.extend(_find_files_in_folder(path))
        else:
            # Single file
            file_paths.append(str(path))
    elif isinstance(path_input, list):
        # List of paths
        for p in path_input:
            path = Path(p)
            if path.is_dir():
                file_paths.extend(_find_files_in_folder(path))
            else:
                file_paths.append(str(path))
    else:
        logger.error(f"Invalid input type: {type(path_input)}")
        return []

    # Inspect each file
    results = []
    for file_path in file_paths:
        metadata = _inspect_single_file(file_path)
        if metadata:
            results.append(metadata)

    return results


def _find_files_in_folder(folder_path: Path) -> List[str]:
    """Recursively find all supported files in a folder."""
    files = []
    try:
        for item in folder_path.rglob("*"):
            if item.is_file() and item.suffix.lower() in SUPPORTED_EXTENSIONS:
                files.append(str(item))
        logger.info(f"Found {len(files)} supported files in {folder_path.name}")
    except Exception as e:
        logger.error(f"Error scanning folder: {e}")
    return files


def _inspect_single_file(file_path: str) -> Optional[Dict[str, Any]]:
    """
    Inspect a single file and return its metadata.

    Returns:
        {
            "file_name": str,
            "file_path": str,
            "file_type": "pdf" | "pptx" | "docx" | "image" | "text",
            "file_size_kb": float,
            "total_pages": int,
            "is_scanned": bool,
            "has_toc": bool,
            "toc_entry_count": int
        }

        Or None if file doesn't exist or is unsupported.
    """
    path = Path(file_path)
    if not path.exists():
        logger.warning(f"File not found: {file_path}")
        return None

    ext = path.suffix.lower()
    size_kb = round(path.stat().st_size / 1024, 1)

    file_type = _ext_to_type(ext)
    if file_type == "unknown":
        logger.warning(f"Unsupported file type: {ext}")
        return None

    metadata = {
        "file_name": path.name,
        "file_path": str(path),
        "file_type": file_type,
        "file_size_kb": size_kb,
        "total_pages": 0,
        "is_scanned": False,
        "has_toc": False,
        "toc_entry_count": 0,
    }

    # Route to type-specific inspector
    if ext == ".pdf":
        _inspect_pdf(metadata)
    elif ext == ".pptx":
        _inspect_pptx(metadata)
    elif ext == ".docx":
        _inspect_docx(metadata)
    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        _inspect_image(metadata)
    elif ext in (".txt", ".md"):
        _inspect_text(metadata)

    logger.info(
        f"Inspected {metadata['file_name']}: {metadata['file_type']}, "
        f"pages={metadata['total_pages']}, scanned={metadata['is_scanned']}, "
        f"toc={metadata['has_toc']}"
    )
    return metadata


# Backward compat: old function name
def inspect_file(file_path: str) -> Optional[Dict[str, Any]]:
    """Deprecated: use inspect() instead. Inspects a single file."""
    results = inspect(file_path)
    return results[0] if results else None


def peek_file(file_path: str, metadata: Dict[str, Any]) -> Optional[str]:
    """
    Sample strategic pages from a file to give a meaningful overview.

    Sampling strategy per type:
    - PDF (has_toc=True)  → pages 1, 2, 3  (intro + early content)
    - PDF (has_toc=False) → pages 1, middle, last  (spread across doc)
    - PPTX                → slides 1-5 full text
    - DOCX                → first 60 paragraphs
    - Text/MD             → first 3000 chars
    - Image               → returns None (caller uses vision)

    Returns full page text from sampled pages (not truncated mid-page).
    """
    file_type = metadata.get("file_type", "")
    total_pages = metadata.get("total_pages", 1)
    has_toc = metadata.get("has_toc", False)

    try:
        if file_type == "pdf":
            import fitz
            doc = fitz.open(file_path)

            # Pick strategic pages (0-indexed internally)
            if has_toc:
                # Book: read first 3 pages — intro + first real content
                pages_to_sample = [i for i in range(min(3, total_pages))]
            else:
                # No TOC: spread across doc — start, middle, end
                mid = total_pages // 2
                pages_to_sample = sorted(set([0, mid, total_pages - 1]))

            # Extract full page text (not truncated mid-page)
            parts = []
            for page_num in pages_to_sample:
                if page_num < total_pages:
                    try:
                        text = doc[page_num].get_text().strip()
                        if text:
                            parts.append(f"[Page {page_num + 1}]\n{text}")
                    except Exception:
                        pass

            doc.close()
            return "\n\n---\n\n".join(parts) if parts else None

        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for slide_num in range(min(5, len(prs.slides))):
                slide = prs.slides[slide_num]
                slide_text = "\n".join(
                    shape.text for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                )
                if slide_text:
                    parts.append(f"[Slide {slide_num + 1}]\n{slide_text}")
            return "\n\n---\n\n".join(parts) if parts else None

        elif file_type == "docx":
            from docx import Document
            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs[:60] if p.text.strip())
            return text if text else None

        elif file_type == "text":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(3000)

        elif file_type == "image":
            # Images require vision — handled by document_content_tool.peek_preview
            return None

    except Exception as e:
        logger.debug(f"peek_file failed for {file_path}: {e}")
        return None

    return None


def _inspect_pdf(metadata: Dict[str, Any]) -> None:
    """Inspect PDF using PyMuPDF."""
    try:
        import fitz
        doc = fitz.open(metadata["file_path"])
        total_pages = len(doc)

        raw_toc = doc.get_toc()
        toc_count = len(raw_toc)
        has_toc = toc_count > 3

        # Sample more pages for better scanned detection
        # Strategy: sample ~10% of document (but at least 3, max 20 pages)
        sample_count = max(3, min(20, total_pages // 10))
        sample_indices = [
            int(total_pages * i / sample_count)
            for i in range(sample_count)
        ]

        total_chars = 0
        for idx in sample_indices:
            if idx < total_pages:
                total_chars += len(doc[idx].get_text())

        avg_chars = total_chars / max(len(sample_indices), 1)
        is_scanned = avg_chars < 300 and not has_toc

        doc.close()

        metadata["total_pages"] = total_pages
        metadata["has_toc"] = has_toc
        metadata["toc_entry_count"] = toc_count
        metadata["is_scanned"] = is_scanned
        metadata["avg_chars_per_page"] = round(avg_chars, 1)  # Include for transparency

    except Exception as e:
        logger.debug(f"PDF inspection failed: {e}")


def _inspect_pptx(metadata: Dict[str, Any]) -> None:
    """Inspect PowerPoint."""
    try:
        from pptx import Presentation
        prs = Presentation(metadata["file_path"])
        total_pages = len(prs.slides)
        metadata["total_pages"] = total_pages
        # PPTX slides always have structure
        metadata["has_toc"] = total_pages > 0
        metadata["toc_entry_count"] = total_pages
    except Exception as e:
        logger.debug(f"PPTX inspection failed: {e}")


def _inspect_docx(metadata: Dict[str, Any]) -> None:
    """Inspect Word document."""
    try:
        from docx import Document
        doc = Document(metadata["file_path"])
        heading_count = sum(
            1 for p in doc.paragraphs
            if p.style and p.style.name.startswith("Heading")
        )
        metadata["has_toc"] = heading_count > 3
        metadata["toc_entry_count"] = heading_count
        metadata["total_pages"] = max(1, len(doc.paragraphs) // 25)
    except Exception as e:
        logger.debug(f"DOCX inspection failed: {e}")


def _inspect_image(metadata: Dict[str, Any]) -> None:
    """Inspect image — always 1 page, marked as scanned."""
    metadata["total_pages"] = 1
    metadata["is_scanned"] = True


def _inspect_text(metadata: Dict[str, Any]) -> None:
    """Inspect text/markdown file."""
    try:
        size = Path(metadata["file_path"]).stat().st_size
        metadata["total_pages"] = max(1, size // 3000)
    except Exception:
        metadata["total_pages"] = 1


def _ext_to_type(ext: str) -> str:
    """Map file extension to type."""
    mapping = {
        ".pdf": "pdf",
        ".pptx": "pptx",
        ".docx": "docx",
        ".png": "image",
        ".jpg": "image",
        ".jpeg": "image",
        ".webp": "image",
        ".txt": "text",
        ".md": "text",
    }
    return mapping.get(ext, "unknown")
