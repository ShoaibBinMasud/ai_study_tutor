"""
PPTX Processor — PowerPoint presentations.

Extracts slide text, speaker notes, and slide titles.
Each slide becomes a section. Speaker notes often have richer content.
"""

import logging
from pathlib import Path
from typing import List

from models.data_models import ProcessedDocument, TOCEntry, Section

logger = logging.getLogger(__name__)


class PPTXProcessor:
    """Processes .pptx files into ProcessedDocument format."""

    def process(self, file_path: str) -> ProcessedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx required: pip install python-pptx")

        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX not found: {file_path}")

        prs = Presentation(file_path)
        label = path.stem.replace("_", " ").replace("-", " ").title()

        toc_entries = []
        sections = []

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            title = self._get_slide_title(slide) or f"Slide {slide_num}"
            body_text = self._get_slide_body(slide)
            notes_text = self._get_slide_notes(slide)

            # Combine body + notes; notes often have more explanation
            content_parts = []
            if body_text:
                content_parts.append(f"[Slide {slide_num}: {title}]\n{body_text}")
            if notes_text:
                content_parts.append(f"[Speaker Notes]\n{notes_text}")

            content = "\n\n".join(content_parts) if content_parts else f"[Slide {slide_num}: {title}]"
            section_id = f"slide_{slide_num}"

            toc_entries.append(TOCEntry(
                title=title,
                level=2,
                page_or_slide=slide_num,
                source_file=file_path,
                section_id=section_id,
            ))
            sections.append(Section(
                section_id=section_id,
                title=title,
                content=content,
                page_start=slide_num,
                page_end=slide_num,
                source_file=file_path,
            ))

        return ProcessedDocument(
            file_path=file_path,
            file_type="pptx",
            label=label,
            toc=toc_entries,
            sections=sections,
        )

    def _get_slide_title(self, slide) -> str:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
        return ""

    def _get_slide_body(self, slide) -> str:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
        return "\n".join(parts)

    def _get_slide_notes(self, slide) -> str:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            return slide.notes_slide.notes_text_frame.text.strip()
        return ""
