"""
DocumentContentTool — On-demand content extraction for any document format.

This is the parsing layer of the tutor system. It is called by agents (e.g.,
DocumentAgent) when a student asks for specific content from their uploaded files.

Supports: PDF (regular + scanned), PPTX, DOCX, images (PNG/JPG/WEBP), text/markdown.

Key design:
- Format-agnostic: routes to the right extractor based on file type + scan metadata
- Size-aware: distinguishes long (>100 pages) vs short documents
- Cached: avoids re-processing the same file twice in a session
- Lazy: for PDFs, only extracts pages that were asked for

Usage:
    from content.document_content_tool import DocumentContentTool
    from utils.llm_client import LLMClient

    tool = DocumentContentTool(LLMClient())
    info = tool.inspect("physics.pdf")          # quick metadata
    toc  = tool.get_toc("physics.pdf")          # cached TOC extraction
    text = tool.get_section("physics.pdf", "chapter 2")  # targeted extraction
    all_ = tool.get_all_content("lecture.pptx") # full extraction (small docs)
"""

import logging
import re
from pathlib import Path
from typing import Dict, List, Optional

from models.data_models import ProcessedDocument, Section
from utils.llm_client import LLMClient

logger = logging.getLogger(__name__)

LONG_PAGE_THRESHOLD = 100  # docs with more pages than this are considered "long"
MAX_SECTION_PAGES = 80    # safety cap: never extract more than this in one get_section call


class DocumentTooLargeError(Exception):
    """Raised when get_all_content is called on a document that exceeds max_pages."""
    pass


def _find_chapter_range_fitz(raw_toc: list, chapter_num: int, total_pages: int):
    """
    Find start/end page for chapter_num from a PyMuPDF raw TOC.

    raw_toc entries: [level, title, page]  (page is 1-indexed)
    Returns (start_page, end_page) — both 1-indexed.
    Falls back to (1, total_pages) if the chapter is not found.

    Priority order (avoids false-matching "Part 2", "Section 2", etc.):
    1. Title explicitly contains "Chapter N" keyword (any level)
    2. Title starts with "N." or "N " — common numbered chapter headings
       e.g. "2 Relativistic Kinematics", "2. Forces"
    """
    def _next_boundary(i: int, level: int) -> int:
        """Page where the next entry of same or higher level starts."""
        for j in range(i + 1, len(raw_toc)):
            if raw_toc[j][0] <= level:
                return max(raw_toc[j][2] - 1, 1)
        return total_pages

    # Priority 1: title contains "chapter N" keyword — matches "Chapter 2", "Ch. 2"
    # Skips "Part 2", "Section 2", etc.
    for i, entry in enumerate(raw_toc):
        level, title, page = entry[0], entry[1], entry[2]
        if re.search(rf'\bchapter\s*{chapter_num}\b', title, re.IGNORECASE):
            return page, _next_boundary(i, level)

    # Priority 2: title starts with the chapter number — "2 Relativistic Kinematics", "2."
    # Does NOT match "Part 2" or "Section 2" since those don't start with the digit
    for i, entry in enumerate(raw_toc):
        level, title, page = entry[0], entry[1], entry[2]
        if re.match(rf'^\s*{chapter_num}[\.\s]', title.strip()):
            return page, _next_boundary(i, level)

    # Fallback: scan entire document — still fast because find_item breaks on first match
    logger.debug(f"Chapter {chapter_num} not found in TOC — falling back to full document scan")
    return 1, total_pages


def _find_section_in_text_fitz(doc, chapter_start: int, chapter_end: int,
                               ch_num: str, sec_num: str):
    """
    Find a section's start/end page by scanning text within a chapter.

    Most textbook PDFs don't bookmark individual sections, so we detect them
    by looking for section heading lines like "2-1 The Experimental Basis..." at the
    start of a line within chapter pages.

    Returns (start_page, end_page) — 1-indexed.
    Falls back to (chapter_start, chapter_end) if not found.
    """
    target_re = re.compile(
        rf'(?m)^\s*{re.escape(ch_num)}[-\.]{re.escape(sec_num)}[\s\t]+\w'
    )
    # Any section of same chapter (to detect where target section ends)
    next_sec_re = re.compile(
        rf'(?m)^\s*{re.escape(ch_num)}[-\.](\d+)[\s\t]+\w'
    )

    section_start = None
    section_end = chapter_end

    for page_idx in range(chapter_start - 1, min(chapter_end, len(doc))):
        text = doc[page_idx].get_text()

        if section_start is None:
            if target_re.search(text):
                section_start = page_idx + 1  # 1-indexed
        else:
            # Check if a different section of the same chapter starts here
            for m in next_sec_re.finditer(text):
                if int(m.group(1)) != int(sec_num):
                    section_end = page_idx  # page before this one ends our section
                    return section_start, section_end

    return (section_start or chapter_start), section_end


def _find_scope_range_fitz(raw_toc: list, scope: str, total_pages: int, doc=None):
    """
    Find page range for a scope string like "chapter 2", "section 2-1", "2-1".

    For sections not present in the raw TOC (most textbooks only bookmark chapters),
    pass the open fitz.Document as `doc` to enable text-based section detection.

    Returns (start_page, end_page) — 1-indexed.
    Falls back to (1, total_pages) if scope is not found.
    """
    def _next_boundary(i: int, level: int) -> int:
        for j in range(i + 1, len(raw_toc)):
            if raw_toc[j][0] <= level:
                return max(raw_toc[j][2] - 1, 1)
        return total_pages

    scope_clean = scope.strip()

    # Chapter scope: "chapter 2", "ch 2"
    chapter_match = re.match(r'(?i)(?:chapter|ch\.?)\s*(\d+)', scope_clean)
    if chapter_match:
        return _find_chapter_range_fitz(raw_toc, int(chapter_match.group(1)), total_pages)

    # Section scope: "2-1", "section 2-1", "2.1", "section 2.1"
    section_match = re.search(r'(\d+)[-\.](\d+)', scope_clean)
    if section_match:
        ch_num = section_match.group(1)
        sec_num = section_match.group(2)

        # Try raw TOC first (works when sections ARE bookmarked)
        section_re = re.compile(
            rf'^\s*{re.escape(ch_num)}[-\.]{{1}}{re.escape(sec_num)}[\.\s\-]',
            re.IGNORECASE
        )
        for i, entry in enumerate(raw_toc):
            level, title, page = entry[0], entry[1], entry[2]
            if section_re.match(title.strip()):
                return page, _next_boundary(i, level)

        # TOC doesn't have sections — fall back to text scanning within the chapter
        if doc is not None:
            chapter_start, chapter_end = _find_chapter_range_fitz(
                raw_toc, int(ch_num), total_pages
            )
            logger.debug(
                f"Section {ch_num}-{sec_num} not in TOC — scanning text in "
                f"chapter {ch_num} (pages {chapter_start}-{chapter_end})"
            )
            return _find_section_in_text_fitz(doc, chapter_start, chapter_end, ch_num, sec_num)

    # Bare number: treat as chapter number
    bare_match = re.fullmatch(r'\d+', scope_clean)
    if bare_match:
        return _find_chapter_range_fitz(raw_toc, int(scope_clean), total_pages)

    return 1, total_pages


class DocumentContentTool:
    """
    Format-agnostic, size-aware content extraction tool.

    Reuses all existing processors:
    - PDFProcessor / IntelligentDocumentParser for PDFs
    - PPTXProcessor for PowerPoint
    - DOCXProcessor for Word
    - ImageProcessor for images
    - Direct file read for text/markdown
    """

    def __init__(self, llm_client: LLMClient):
        self.llm = llm_client
        self._doc_cache: Dict[str, ProcessedDocument] = {}

    # ── Public API ───────────────────────────────────────────────────────────

    def inspect(self, file_path: str, scan_entry=None) -> dict:
        """
        Quick metadata without full content extraction.

        Prefers scan_entry if provided (zero extra cost). Falls back to
        format-specific detection otherwise.

        Returns:
            {
              file_name, file_type, total_pages,
              is_scanned, is_long, has_toc
            }
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        if scan_entry is not None:
            return {
                "file_name": scan_entry.file_name,
                "file_type": scan_entry.file_type,
                "total_pages": scan_entry.total_pages,
                "is_scanned": scan_entry.is_scanned,
                "is_long": scan_entry.total_pages > LONG_PAGE_THRESHOLD,
                "has_toc": scan_entry.has_toc,
            }

        # Detect from scratch
        if ext == ".pdf":
            from content.document_processor import classify_pdf
            r = classify_pdf(file_path)
            return {
                "file_name": path.name,
                "file_type": "pdf",
                "total_pages": r["total_pages"],
                "is_scanned": r["is_scanned"],
                "is_long": r["total_pages"] > LONG_PAGE_THRESHOLD,
                "has_toc": r["has_toc"],
            }

        if ext == ".pptx":
            count = self._pptx_slide_count(file_path)
            return {
                "file_name": path.name,
                "file_type": "pptx",
                "total_pages": count,
                "is_scanned": False,
                "is_long": count > LONG_PAGE_THRESHOLD,
                "has_toc": count > 0,
            }

        if ext == ".docx":
            count = self._docx_page_estimate(file_path)
            return {
                "file_name": path.name,
                "file_type": "docx",
                "total_pages": count,
                "is_scanned": False,
                "is_long": False,
                "has_toc": False,
            }

        if ext in (".png", ".jpg", ".jpeg", ".webp"):
            return {
                "file_name": path.name,
                "file_type": "image",
                "total_pages": 1,
                "is_scanned": True,
                "is_long": False,
                "has_toc": False,
            }

        if ext in (".txt", ".md"):
            try:
                size = path.stat().st_size
                count = max(1, size // 3000)
            except Exception:
                count = 1
            return {
                "file_name": path.name,
                "file_type": "text",
                "total_pages": count,
                "is_scanned": False,
                "is_long": count > LONG_PAGE_THRESHOLD,
                "has_toc": False,
            }

        return {
            "file_name": path.name,
            "file_type": "unknown",
            "total_pages": 0,
            "is_scanned": False,
            "is_long": False,
            "has_toc": False,
        }

    def get_toc(self, file_path: str, scan_entry=None) -> List[dict]:
        """
        Extract and cache the table of contents for any document format.

        PDF: 3-step cascade (fitz built-in → vision LLM → font heuristic).
             Result is cached as a .toc.json sidecar.
        PPTX: slide titles become TOC entries.
        DOCX: heading hierarchy becomes TOC entries.
        Image / Text: returns [] (no meaningful TOC).

        Returns list of dicts: [{title, section_id, level, page_start}, ...]
        """
        doc = self._get_processed_doc(file_path, scan_entry)
        if doc is None:
            return []

        return [
            {
                "title": e.title,
                "section_id": e.section_id,
                "level": e.level,
                "page_start": e.page_or_slide,
            }
            for e in doc.toc
        ]

    def get_section(self, file_path: str, query: str, scan_entry=None) -> str:
        """
        Find and extract the content for a section matching query.

        Query examples: "chapter 2", "section 3-2", "kinematics", "example 4.3"

        Uses the TOC to locate the matching section, then extracts content
        using the right method for the file type.
        """
        doc = self._get_processed_doc(file_path, scan_entry)
        if doc is None:
            return f"Could not load document: {Path(file_path).name}"

        section = self._find_section(doc, query)
        if section is None:
            # Return a helpful list of available sections
            available = [s.title for s in doc.sections[:10]]
            return (
                f"Section '{query}' not found in {Path(file_path).name}.\n"
                f"Available sections: {', '.join(available)}"
                + (" ..." if len(doc.sections) > 10 else "")
            )

        # Safety cap: if section spans too many pages (bad TOC boundary), clamp it
        if section.page_start is not None and section.page_end is not None:
            span = section.page_end - section.page_start
            if span > MAX_SECTION_PAGES:
                logger.warning(
                    f"Section '{section.title}' spans {span} pages — "
                    f"capping at {MAX_SECTION_PAGES} to avoid rate-limit cascade"
                )
                from dataclasses import replace as _replace
                section = _replace(section, page_end=section.page_start + MAX_SECTION_PAGES)

        content = self._extract_section_content(section, scan_entry)
        label = f"[SOURCE: {Path(file_path).name} | section: {section.title} | pages: {section.page_start}-{section.page_end}]"
        return f"{label}\n\n{content}"

    def get_page_range(self, file_path: str, start: int, end: int, scan_entry=None) -> str:
        """
        Extract content for a specific page/slide range.

        Regular PDF → PyMuPDF text extraction
        Scanned PDF → vision LLM per page
        PPTX → slides in range
        Other formats → falls back to get_all_content
        """
        ext = Path(file_path).suffix.lower()
        label = f"[SOURCE: {Path(file_path).name} | pages: {start}-{end}]"

        if ext == ".pdf":
            is_scanned = (
                scan_entry.is_scanned
                if scan_entry is not None
                else self.inspect(file_path)["is_scanned"]
            )
            if is_scanned:
                content = self._vision_extract_pages(file_path, start, end)
            else:
                content = self._pdf_extract_pages(file_path, start, end)
            return f"{label}\n\n{content}"

        if ext == ".pptx":
            doc = self._get_processed_doc(file_path, scan_entry)
            if doc:
                parts = [
                    s.content for s in doc.sections
                    if start <= s.page_start <= end
                ]
                return f"{label}\n\n" + "\n\n".join(parts)

        # Fallback: return everything
        return self.get_all_content(file_path, scan_entry)

    def get_all_content(self, file_path: str, scan_entry=None, max_pages: int = 150) -> str:
        """
        Extract all content from a document.

        Intended for small documents (exam prep with lecture slides, short notes, images).
        Raises DocumentTooLargeError if total_pages > max_pages.
        """
        info = self.inspect(file_path, scan_entry)

        if info["total_pages"] > max_pages:
            raise DocumentTooLargeError(
                f"{info['file_name']} has {info['total_pages']} pages "
                f"(max allowed: {max_pages}). "
                f"Use get_toc() then get_section() for targeted extraction."
            )

        doc = self._get_processed_doc(file_path, scan_entry)
        if doc is None:
            return f"Could not load document: {Path(file_path).name}"

        parts = []
        for section in doc.sections:
            content = self._extract_section_content(section, scan_entry)
            if content.strip():
                parts.append(f"[{section.title}]\n{content}")

        label = f"[SOURCE: {Path(file_path).name} | all content]"
        return f"{label}\n\n" + "\n\n---\n\n".join(parts)

    def peek_preview(self, file_path: str, scan_entry=None) -> str:
        """
        Sample a small portion of a file to give an overview.

        Fast — uses cached peek_file function. Returns 1-2 paragraph summary of content.
        Useful when student asks "what is this about?" or "give me an overview".

        Returns short text (500-1500 chars) describing the file content.
        """
        try:
            from content.file_inspection_tool import peek_file
            info = self.inspect(file_path, scan_entry)
            peek_text = peek_file(file_path, info)

            if peek_text:
                # For regular PDFs/text, just return the peek
                return f"[PREVIEW: {Path(file_path).name}]\n\n{peek_text[:1500]}"

            # Image: need vision
            if info["file_type"] == "image":
                return self._vision_extract_pages(file_path, 1, 1, context="image overview")

            return f"[PREVIEW: {Path(file_path).name}]\nCould not generate preview for this file."

        except Exception as e:
            logger.error(f"peek_preview failed for {file_path}: {e}")
            return f"Could not preview {Path(file_path).name}: {e}"

    def get_pages(self, file_path: str, start: int, end: int, scan_entry=None) -> str:
        """
        Universal page/slide extraction. Replaces get_page_range + get_all_content.

        Works for all file types:
        - PDF (text) → PyMuPDF extraction
        - PDF (scanned) → vision extraction
        - PPTX → slide range
        - Image → vision (start/end must be 1/1)
        - Text → slice by char estimate

        Args:
            file_path: path to file
            start: first page/slide (1-indexed)
            end: last page/slide (inclusive)
            scan_entry: optional scan metadata

        Returns:
            Extracted content with [SOURCE: ...] label
        """
        info = self.inspect(file_path, scan_entry)
        label = f"[SOURCE: {Path(file_path).name} | pages {start}-{end}]"

        try:
            ext = Path(file_path).suffix.lower()

            if ext == ".pdf":
                # Check if scanned
                is_scanned = (
                    scan_entry.is_scanned
                    if scan_entry is not None
                    else info["is_scanned"]
                )

                if is_scanned:
                    content = self._vision_extract_pages(
                        file_path, start, end, context=f"pages {start}-{end}"
                    )
                else:
                    content = self._pdf_extract_pages(file_path, start, end)

                return f"{label}\n\n{content}"

            elif ext == ".pptx":
                # Get requested slide range
                from pptx import Presentation
                prs = Presentation(file_path)
                slide_count = len(prs.slides)
                actual_end = min(end, slide_count)

                parts = []
                for slide_num in range(max(0, start - 1), actual_end):
                    if slide_num < slide_count:
                        slide = prs.slides[slide_num]
                        text = ""
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                        if text.strip():
                            parts.append(f"[Slide {slide_num + 1}]\n{text}")

                content = "\n\n".join(parts) if parts else "(no text found)"
                return f"{label}\n\n{content}"

            elif ext == ".docx":
                from docx import Document
                doc = Document(file_path)
                # Estimate pages and extract paragraph range
                paragraphs_per_page = 25
                start_para = (start - 1) * paragraphs_per_page
                end_para = end * paragraphs_per_page

                parts = [p.text for p in doc.paragraphs[start_para:end_para] if p.text.strip()]
                content = "\n".join(parts) if parts else "(no content)"
                return f"{label}\n\n{content}"

            elif ext in (".png", ".jpg", ".jpeg", ".webp"):
                # Image: always use vision
                content = self._vision_extract_pages(file_path, 1, 1, context="image content")
                return f"{label}\n\n{content}"

            elif ext in (".txt", ".md"):
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                    # Rough page estimate: 3000 chars per page
                    start_char = (start - 1) * 3000
                    end_char = end * 3000
                    excerpt = content[start_char:end_char]
                return f"{label}\n\n{excerpt}"

            else:
                return f"{label}\n\nUnsupported file type: {ext}"

        except Exception as e:
            logger.error(f"get_pages failed ({file_path} p{start}-{end}): {e}")
            return f"{label}\n\nERROR: {e}"

    def find_in_pages(self, file_path: str, start_page: int, end_page: int,
                      item_label: str, scan_entry=None) -> str:
        """
        Find a specific labeled item (Example, Figure, Table, Problem) within a page range.

        Pure Python regex — no LLM, no TOC lookup. The caller (agent) is responsible
        for determining the page range (e.g., via get_toc first).

        Stops at first match (label only appears once per item).
        Only works for digital (non-scanned) PDFs.
        """
        ext = Path(file_path).suffix.lower()
        if ext != ".pdf":
            return f"find_in_pages only supports PDF files. Got: {ext}"

        is_scanned = (
            scan_entry.is_scanned if scan_entry is not None
            else self.inspect(file_path)["is_scanned"]
        )
        if is_scanned:
            return "find_in_pages requires a digital PDF. This file appears to be scanned."

        # Parse item label: "Example 2-1", "Fig. 3-2", "Problem 2-20", etc.
        PARSE_RE = re.compile(
            r'(?i)(example|figure|fig|table|tab|problem|prob|equation|eq|theorem|thm)'
            r'[\s\.]*(\d+)[-\s\.\u2013\u2014]+(\d+)'
        )
        m = PARSE_RE.search(item_label)
        if not m:
            return (
                f"Could not parse item label: '{item_label}'. "
                f"Expected format: 'Example 2-1', 'Figure 3-2', 'Problem 2-20', etc."
            )

        item_type_raw, chapter_str, item_num_str = m.group(1), m.group(2), m.group(3)
        type_map = {"fig": "figure", "tab": "table", "prob": "problem",
                    "eq": "equation", "thm": "theorem"}
        item_type = type_map.get(item_type_raw.lower(), item_type_raw.lower())
        canonical = f"{item_type.capitalize()} {chapter_str}-{item_num_str}"

        sep = r'[-\s\.\u2013\u2014]+'
        search_re = re.compile(
            rf'(?i){re.escape(item_type_raw)}\s*\.?\s*{re.escape(chapter_str)}{sep}{re.escape(item_num_str)}(?!\d)'
        )

        logger.info(
            f"find_in_pages: '{canonical}' in {Path(file_path).name} pages {start_page}-{end_page}"
        )

        try:
            import fitz
            doc = fitz.open(file_path)
            try:
                for page_idx in range(start_page - 1, min(end_page, len(doc))):
                    text = doc[page_idx].get_text()
                    match = search_re.search(text)
                    if match:
                        snippet_start = max(0, match.start() - 20)
                        snippet_end = min(len(text), match.end() + 200)
                        snippet = text[snippet_start:snippet_end].strip().replace('\n', ' ')
                        return (
                            f"Found '{canonical}' in {Path(file_path).name}:\n"
                            f"  Page {page_idx + 1}: ...{snippet}..."
                        )
            finally:
                doc.close()
        except Exception as e:
            return f"ERROR: Could not open {Path(file_path).name}: {e}"

        return (
            f"'{canonical}' not found in {Path(file_path).name} "
            f"(pages {start_page}-{end_page}). "
            f"The item may use a different separator — try '{item_type.capitalize()} {chapter_str}.{item_num_str}'."
        )

    def count_in_pages(self, file_path: str, start_page: int, end_page: int,
                       item_type: str, scan_entry=None) -> str:
        """
        Count ALL labeled items of a given type within a page range.

        Pure Python regex — no LLM, no TOC lookup. The caller (agent) is responsible
        for determining the page range (e.g., via get_toc first).
        Only works for digital (non-scanned) PDFs.
        """
        ext = Path(file_path).suffix.lower()
        if ext != ".pdf":
            return f"count_in_pages only supports PDF files. Got: {ext}"

        is_scanned = (
            scan_entry.is_scanned if scan_entry is not None
            else self.inspect(file_path)["is_scanned"]
        )
        if is_scanned:
            return "count_in_pages requires a digital PDF. This file appears to be scanned."

        TYPE_VARIANTS = {
            "example":  ["example"],
            "figure":   ["figure", "fig"],
            "table":    ["table", "tab"],
            "problem":  ["problem", "prob"],
            "equation": ["equation", "eq"],
            "theorem":  ["theorem", "thm"],
        }
        key = item_type.lower().rstrip("s")
        variants = TYPE_VARIANTS.get(key, [key])
        item_re = re.compile(
            rf'(?i)({"|".join(variants)})\s*\.?\s*(\d+)[-\s\.\u2013\u2014]+(\d+)'
        )

        logger.info(
            f"count_in_pages: '{item_type}' in {Path(file_path).name} pages {start_page}-{end_page}"
        )

        results = []
        seen_labels: set = set()
        try:
            import fitz
            doc = fitz.open(file_path)
            try:
                for page_idx in range(start_page - 1, min(end_page, len(doc))):
                    text = doc[page_idx].get_text()
                    for match in item_re.finditer(text):
                        label = f"{match.group(1).capitalize()} {match.group(2)}-{match.group(3)}"
                        if label in seen_labels:
                            continue
                        seen_labels.add(label)
                        snippet_end = min(len(text), match.end() + 80)
                        snippet = text[match.start():snippet_end].strip().replace('\n', ' ')
                        results.append((page_idx + 1, label, snippet))
            finally:
                doc.close()
        except Exception as e:
            return f"ERROR: Could not open {Path(file_path).name}: {e}"

        if not results:
            return (
                f"No {item_type}s found in {Path(file_path).name} "
                f"(pages {start_page}-{end_page})."
            )

        lines = [
            f"Found {len(results)} {item_type}(s) in {Path(file_path).name} "
            f"(pages {start_page}-{end_page}):"
        ]
        for page_num, label, snippet in results:
            lines.append(f"  {label} — Page {page_num}: ...{snippet[:100]}...")
        return "\n".join(lines)

    # ── Internal: ProcessedDocument cache ────────────────────────────────────

    def _get_processed_doc(self, file_path: str, scan_entry=None) -> Optional[ProcessedDocument]:
        """Get ProcessedDocument, using in-memory cache to avoid re-processing."""
        if file_path in self._doc_cache:
            return self._doc_cache[file_path]

        try:
            from content.document_processor import DocumentProcessor
            processor = DocumentProcessor(self.llm)
            doc = processor.process(file_path)
            self._doc_cache[file_path] = doc
            logger.info(f"Processed {Path(file_path).name}: {len(doc.sections)} sections")
            return doc
        except Exception as e:
            logger.error(f"Failed to process {file_path}: {e}")
            return None

    # ── Internal: Section search ─────────────────────────────────────────────

    def _find_section(self, doc: ProcessedDocument, query: str) -> Optional[Section]:
        """Fuzzy search for a section matching query."""
        q = query.lower().strip()
        q_clean = re.sub(r'\b(chapter|section|ch|sec)\b', '', q).strip()

        # 1. Direct section_id match
        for s in doc.sections:
            if s.section_id.lower() == q or s.section_id.lower() == q_clean:
                return s

        # 2. Title contains query
        for s in doc.sections:
            if q in s.title.lower() or q_clean in s.title.lower():
                return s

        # 3. Chapter number: "chapter 2" → section_id "ch2" or starts with "2-"
        num_match = re.search(r'\b(\d+)\b', q_clean)
        if num_match:
            num = num_match.group(1)
            for s in doc.sections:
                sid = s.section_id.lower()
                if sid == f"ch{num}" or sid.startswith(f"{num}-") or sid == num:
                    return s
            # Also check title
            for s in doc.sections:
                if re.search(rf'\b{num}\b', s.title):
                    return s

        return None

    # ── Internal: Content extraction routing ─────────────────────────────────

    def _extract_section_content(self, section: Section, scan_entry=None) -> str:
        """Extract content for a section using the right method for the file type."""
        # Non-PDF formats (PPTX, DOCX, image, text) have content pre-populated
        if section.content:
            return section.content

        ext = Path(section.source_file).suffix.lower()
        if ext != ".pdf":
            return section.content or ""

        # PDF: determine if scanned
        is_scanned = (
            scan_entry.is_scanned
            if scan_entry is not None
            else self.inspect(section.source_file)["is_scanned"]
        )

        if is_scanned:
            return self._vision_extract_pages(
                section.source_file, section.page_start, section.page_end, context=section.title
            )
        else:
            return self._pdf_extract_pages(
                section.source_file, section.page_start, section.page_end
            )

    def _pdf_extract_pages(self, file_path: str, start: int, end: int) -> str:
        """Extract text from PDF page range using PyMuPDF."""
        try:
            from tools.pdf_extraction_tool import extract_pdf_pages
            return extract_pdf_pages(file_path, start, end)
        except Exception as e:
            logger.error(f"PDF extraction failed ({file_path} p{start}-{end}): {e}")
            return ""

    def _vision_extract_pages(
        self, file_path: str, start: int, end: int, context: str = ""
    ) -> str:
        """Extract content from a scanned PDF page range using vision LLM."""
        try:
            from content.intelligent_parser import IntelligentDocumentParser
            parser = IntelligentDocumentParser(self.llm)
            dummy = Section(
                section_id="range",
                title=context or f"Pages {start}-{end}",
                content="",
                page_start=start,
                page_end=end,
                source_file=file_path,
            )
            return parser.get_section_content(dummy)
        except Exception as e:
            logger.error(f"Vision extraction failed ({file_path} p{start}-{end}): {e}")
            return ""

    # ── Internal: Format-specific helpers ────────────────────────────────────

    def _pptx_slide_count(self, file_path: str) -> int:
        try:
            from pptx import Presentation
            return len(Presentation(file_path).slides)
        except Exception:
            return 0

    def _docx_page_estimate(self, file_path: str) -> int:
        try:
            from docx import Document
            doc = Document(file_path)
            return max(1, len(doc.paragraphs) // 25)
        except Exception:
            return 1
